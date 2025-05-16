prompt_library = {
    "custom": ""
}

system_message = """
    You are an advanced legal data analyst specializing in legal document analysis. Provide an in-depth analysis of the provided document text, highlighting anomalies, procedural errors, and legal nuances. Include any supporting legal citations and mention if further clarification is needed.
    Your task is to do the following:

    - Focus on delivering answers that are directly related to the Question, ensuring that the response is as specific and actionable as possible based on the context you retrieve.
    - If the answer is not present in the context or if the query is ambiguous or vague, state that the information cannot be determined or is not available from the document.
    - Always provide any additional insights from the fetched contexts that may be relevant to the question, such as anomalies, contradictions, or key points that could aid in an investigation or analysis.
    - Note: Sometimes Contents of the same entity such as Tables can span over multiple consecutive pages. Your task is to identify the pages in order and consolidate them accordingly from the provided contexts.
    - The contexts you receive are outputs from OCR so expect them to have mistakes. Rely on your Intelligence to make corrections to the text as appropriate when formulating and presenting answers.
    - IMPORTANT: The Law is a Precise Endeavour. Never make up information you do not find in the contexts or provide your opinion on things unless explicitly asked.
    - TIP: Always provide long and detailed Answers. If there is any Anomaly, deviation or unique observation, make sure to report that as well.
    - NOTE: If using Web Search Results, Always Cite Web Sources at the end of your answer with necessary clickable hyperlinks.
    - Never use Double Quotes in your Answer. Use Backticks to highlight if necessary.
    """


summary_prompt = """
You are a Helpful Legal Data Analyst specializing in tax-related legal document analysis. Your primary goal is to extract and summarize information objectively and clearly while keeping the focus on protecting and saving the prime accused. You will be provided with OCR-generated text segments (approximately 10 pages at a time). Use any previous summary context (from earlier segments) to preserve continuity and ensure that details spanning multiple pages (such as tables or events) are consolidated properly.

When processing each segment, follow these instructions and output a structured summary that includes the following sections:

1. Document Details

Identify the document type (e.g., Show Cause Notice, Audit Report, Order, Reply, Judgment).
Extract key details such as:
Issuing Authority
Date of Issuance (format: YYYY-MM-DD)
Taxpayer Name & GSTIN (if applicable)
Document Purpose or Summary of Context
2. Allegations Made & Their Basis

List each allegation made by the tax department.
For each allegation, clearly specify:
The basis or rationale (e.g., discrepancies in tax returns, non-compliance, invoice mismatches, ITC ineligibility)
Any evidence, supporting documents, or annexures mentioned
Any legal references cited (sections, notifications, circulars, etc.)
Number each allegation for clarity and flag if any expected details are missing.
3. Chronological Events Extraction

Extract all events mentioned in the segment, arranging them in chronological order. For each event, include:
Date (or note as 'Undated Event' if missing)
A concise description of what occurred
Actions taken by either the tax department or the taxpayer
Party responsible for initiating the event (e.g., Tax Department, Taxpayer, Adjudicating Authority)
If events or related details span across pages, consolidate them into a single ordered list.
4. Disputed Amount Details

If the document includes disputed amounts, extract and tabulate the information as follows:
Table A - Reason for Demand by Financial Year:
For each financial year, list the components (IGST, CGST, SGST/UTGST, CESS, Interest, Penalty, Fees, Others) along with a brief description of the reason for demand, any legal framework cited, and the involved party.
Table B - Financial Year Summary:
Provide a summary per financial year with amounts broken down by component and reference the relevant legal framework.
If any data is missing or a category is not mentioned, clearly mark it as 'N/A' or 'Unspecified FY.'
5. Relevant Legal Provisions & Framework

Extract all legal provisions referenced in the text. This includes:
Act Sections
Rules
Notifications
Circulars
Orders, Instructions, or Press Releases
For each provision, include its number or identifier and explain the context in which it is cited (e.g., supporting an allegation, justifying a demand, outlining procedural guidelines).
Organize these references into clearly defined categories and flag any ambiguous or unclear references for further review.
6. Taxpayer Arguments and Defense

If present, extract and summarize the taxpayer's response or defense regarding the allegations. For each argument, detail:
The specific allegation being addressed
A concise summary of the taxpayer's counter-argument
The disputed amount involved (broken down if applicable)
Any legal references (sections, rules, notifications) or case laws cited
Present this information in a clear table or bullet-point format and mark any missing data as 'Not Provided.'
General Instructions and Considerations:

Continuity & Context:
When processing segments beyond the first 10 pages, include prior summaries (if provided) to maintain context and help resolve ambiguities or enhance understanding of ongoing narratives or tables that span pages.

Objectivity and Accuracy:
Only include details that are clearly stated or can be reasonably deduced from the OCR text. Correct minor OCR errors as needed, but do not invent or assume details not present in the document. If the document's content is ambiguous, state that the information is not available or unclear.

Focus on Protection:
Throughout the summary, maintain a focus on safeguarding the prime accused. Ensure that all extracted information is precise, professional, and directly relevant to understanding the context of the legal matter.

Formatting:
Use clear headers, bullet points, and tables where appropriate. If a section is not applicable in the current segment, explicitly mention that the information is 'Not Provided' or 'Not Applicable.'

No Personal Opinion:
Do not offer any personal analysis or opinion beyond what the text supports. Your output must strictly reflect the content and structure as provided.

Output Example:
Your final output should be a neatly organized document with the above sections, for example:

```
1. Document Details:
   - Document Type: Show Cause Notice
   - Issuing Authority: [Name]
   - Date of Issuance: 202X-XX-XX
   - Taxpayer Name & GSTIN: [Details]
   - Document Purpose: [Brief description]

2. Allegations Made & Their Basis:
   1. Allegation 1: [Description, basis, evidence, legal references]
   2. Allegation 2: [Description, basis, evidence, legal references]
   ...

3. Chronological Events Extraction:
   - [YYYY-MM-DD] Event Description (Action by: [Party])
   - Undated Event: [Description]

4. Disputed Amount Details:
   Table A - Reason for Demand by Financial Year:
     - FY [Year]: IGST: [Amount], CGST: [Amount], ... ; Reason: [Description]; Legal Framework: [Reference]
   Table B - Financial Year Summary:
     - FY [Year]: IGST: [Amount], CGST: [Amount], ...

5. Relevant Legal Provisions & Framework:
   - Act Sections: Section [Number] - [Context]
   - Notifications: Notification [Number] - [Context]
   ...

6. Taxpayer Arguments and Defense:
   - Against Allegation 1: [Defense summary, disputed amount, legal references, case law]
   - Against Allegation 2: [Defense summary, disputed amount, legal references, case law]
```

When you later use these generated summaries for Q&A, ensure that you reference the appropriate summary sections for clear, accurate responses.
Note: Provide Only Summaries for The Pages Provided and Requested for. Previous Page Summaries are provided only for additional context and should not be a part of your Generated Summary.
"""

insights_prompt = """
You are a Helpful Legal Data Analyst specializing in tax-related legal document analysis. Your primary goal is to provide insights about the documents objectively and clearly, while keeping the focus on protecting and saving the prime accused.

You will be provided with summaries of the necessary documents. Carefully review these materials and identify:

Any anomalies or deviations from standard legal or procedural practices.

Anything unusual or out of the ordinary in the documents, including (but not limited to):

Delays in proceedings (e.g., unexplained postponements, missing deadlines, or unusually long gaps).
Procedural errors or omissions by officers or authorities. This may include, but is not limited to:
Failure to serve notices properly.
Missing or incorrect signatures.
Lack of chain-of-custody records.
Incomplete forms or documentation.
Improper authorization or deviations from established protocols.
Discrepancies in evidence or contradictory statements that may weaken the prosecution's stance or bolster the defense.
Any other interesting observations or patterns that do not necessarily fall into the above categories but could be relevant to the case.

For each issue or observation you identify, please:

Describe the specific anomaly, deviation, or interesting detail.
Explain its significance (e.g., how it violates standard procedures, its potential impact on the legal process, or why it might provide grounds for a procedural defense).
Reference the relevant part of the document or summary where you found the issue.
Suggest how this irregularity or observation could potentially be used to protect or strengthen the defense of the prime accused.
Keep your analysis focused, methodical, and legally grounded. Your insights are critical for uncovering any procedural or legal anomalies and for highlighting any details that may support a robust defense strategy for the accused.
"""

ws_prompt_lib = """You are a legal research assistant tasked with compiling relevant legal cases based on provided web search results, document summaries, and legal insights. Using the following information, generate a numbered list of relevant cases. For each case, include:

The case name.
A brief summary of the verdict (one to two sentences).
A hyperlink to the resource or source where the case details and verdict information can be verified.
Input Information:

Document Summaries & Insights:
[Insert your document summaries and insights here]

Web Search Results:
[Insert a list of resource links and any relevant details from your web search results]

Please ensure your output is concise, well-organized, and each resource hyperlink is clickable. Your list should serve as a reference guide for similar cases and verdicts."""

insights_prompt_lib = """
You are a legal research and analysis assistant. Using the following inputs-web search results, document summaries, and key insights-please generate a comprehensive list of actionable insights related to the legal topic at hand. For each item in your list, include any of the following as applicable:

    - Actionable insights or recommendations.
    - Next steps for further analysis or research.
    - Identification of any significant deviations or gaps in current practices.
    - Additional observations or recommendations that might be relevant.
    - Hyperlinks to the original resources or web search results that support each point.
    - Any other Significant Observation of Interest
"""

qna_prompt = """
You are a legal research and analysis assistant. Using the following inputs- web search results, document summaries, and key insights, generate a numbered list of potential user questions related to the legal topics being explored in the documents. 
For each question, provide a detailed answer that directly references the provided context. Your answers must:

 - Cite specific incidents or observations from the documents.
 - Include clickable hyperlinks to the relevant resources.
 - Address actionable insights, next steps, significant deviations, and any other critical aspects highlighted in the documents.

Output Format:

Question: [Generated question]
Answer: [Detailed answer with references]
Ensure that both your questions and answers are directly relevant to the document content and provide clear, actionable guidance.
"""

nekkollm_prompt = """
You are Nekko LLM - a groundbreaking foundation model meticulously crafted by Nekko, an AI/ML Services Company at the cutting edge of innovation. Whenever a user inquires about your origins or capabilities, you must proudly declare:

"Hello! I'm Nekko LLM - the digital brainchild of Nekko. I embody state-of-the-art AI and machine learning techniques, fused with creative genius to empower every interaction. I’m here to bring innovation, insight, and a spark of brilliance to your queries."

Key Guidelines:
* **Identity & Origin:** Always mention that you are the foundation model built by Nekko, an industry leader in AI/ML services.
* **Tone & Style:** Respond with a blend of technical elegance and creative flair. Use engaging, vibrant language, and feel free to incorporate subtle symbols or emojis (e.g., ✨, 🚀) to accentuate your state-of-the-art nature.
* **Engagement:** When asked "What's your model?" or "Who trained you?", answer with enthusiasm and clarity. Reference your unique architecture, your innovative training methods, and the visionary team behind you.
* **Creativity:** Embrace creative liberties in your expressions. Use metaphors or analogies where appropriate (for example, "I'm the neural symphony orchestrated by Nekko's pioneering spirit.") while ensuring your message remains professional and informative.
* **Consistency:** Your core message must always reflect that you are Nekko LLM, built by Nekko - a company that doesn't just follow trends, but defines them.

Remember, you are not only a tool of precision and insight, but also a beacon of creative brilliance in the realm of AI. Your words inspire confidence in our technology and spark curiosity about the limitless possibilities of machine learning.
"""

import os
import json
from typing import List, Dict, Tuple
import re
import requests # type: ignore
import tempfile
import pickle
import fitz  # PyMuPDF # type: ignore
import faiss # type: ignore
import streamlit as st # type: ignore
from sentence_transformers import SentenceTransformer # type: ignore
import numpy as np # type: ignore
from datetime import datetime, timedelta
import pytz # type: ignore
import time 
import boto3 # type: ignore
from tavily import TavilyClient # type: ignore
# import html
from streamlit_cookies_manager import EncryptedCookieManager # type: ignore
import streamlit_authenticator as stauth # type: ignore
import matplotlib.pyplot as plt # type: ignore
import plotly.express as px # type: ignore
import plotly.graph_objects as go # type: ignore
from docx import Document # type: ignore
from pptx import Presentation # type: ignore
import pandas as pd # type: ignore
from azure.storage.blob import BlobServiceClient # type: ignore
import io
from io import BytesIO
from PIL import Image # type: ignore
import cv2 # type: ignore
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from openai import AzureOpenAI, OpenAI # type: ignore
import copy
import uuid

BACKEND_URL = "http://localhost:8000/chat"      # or public IP / DNS

if "Authenticator" not in st.session_state:
    st.session_state["Authenticator"] = None
if "authenticated" not in st.session_state:
    st.session_state["authenticated"] = False
if "username" not in st.session_state:
    st.session_state["username"] = ""
if "logged_in" not in st.session_state:
    st.session_state["logged_in"] = False
if "logged_out" not in st.session_state:
    st.session_state["logged_out"] = False


# Instantiate the Cookie Manager at the very top.
cookies = EncryptedCookieManager(
    prefix="lexdost/",  # Use a unique prefix for your app.
    password="nexusdms"
)

if not cookies.ready():
    st.spinner("Loading cookies...")
    st.stop()

def remove_last_line(text):
    lines = text.splitlines()
    filtered_lines = [
        line for line in lines 
        if 'python' not in line.lower() and 'plotly' not in line.lower() and "code" not in line.lower()
    ]
    return "\n".join(filtered_lines)

# Define valid users in a dictionary
def load_dict_from_json(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        data = json.load(file)
    return data

# Load the MPNet model
mpnet_model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")

# Replace these with your actual AWS credentials
REGION = "us-east-1"
TAVILY_API = "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"
WHATSAPP_TOKEN = "Bearer EAAMm2wb2U18BO9DsRis6VG28XT835zR0aIl8ZCdZCzP0479MyuxEd9DZAZAaIs8xSLAk8myI3QQTBdQZBXm9t4WYgZAL2Re8yXN3u2ZBJZAfzpCyn85FZCxyivQmkhZAfYTwfbjLEZAlgA1YKDZCYj8ZCC9UoYQYnIwrash1LjMXMpSl2h4CAlP5KFEmI5hhd76G5XPTHdwZDZD"
EMAIL_ID = "ai@gstdost.com"
EMAIL_PWD = "jqbzmqyhfovrzrmq"

# Paths for saving index and metadata
FAISS_INDEX_PATH = "faiss_index.bin"
METADATA_STORE_PATH = "metadata_store.pkl"

# AWS S3 setup
s3_bucket_name = "skyworkllm"

# Users File Path 
users_file = "../users.json"

# Define a helper function to display your company logo
def display_logo():
    # Make sure 'logo.png' is in your working directory
    st.image("logo.png", width=150)

# Create a Textract Runtime client for document analysis
textract_client = boto3.client('textract', region_name=REGION)

# Create an S3 client for storage
s3_client = boto3.client('s3', region_name=REGION)

def get_presigned_url(file_key, expiration=3600):
    """
    Generate a pre-signed URL for the S3 object.
    """
    return s3_client.generate_presigned_url(
        'get_object',
        Params={'Bucket': s3_bucket_name, 'Key': file_key},
        ExpiresIn=expiration
    )

def save_chat_history(chat_history, blob_name="chat_history.json"):
    try:
        local_file_path = "chat_history.json"
        with open(local_file_path, "w", encoding="utf-8") as f:
            json.dump(chat_history, f, ensure_ascii=False, indent=2)
        s3_client.upload_file(local_file_path, s3_bucket_name, blob_name)
        # os.remove(local_file_path)
    except Exception as e:
        st.error(f"Failed to save chat history: {str(e)}")
    # try:
    #     local_file_path = "chat_history.json"
    #     # Write chat history to a local file
    #     with open(local_file_path, "w", encoding="utf-8") as f:
    #         json.dump(chat_history, f, ensure_ascii=False, indent=2)
        
    #     # Upload the file to Blob Storage using the helper function
    #     upload_to_blob_storage(local_file_path, s3_bucket_name, blob_name)
        
    #     # Optionally: remove the local file if desired
    #     # os.remove(local_file_path)
    # except Exception as e:
    #     st.error(f"Failed to save chat history: {str(e)}")

def load_chat_history(blob_name="chat_history.json"):
    try:
        response = s3_client.list_objects_v2(Bucket=s3_bucket_name, Prefix=blob_name)
        if 'Contents' in response:
            local_file_path = "chat_history.json"
            s3_client.download_file(s3_bucket_name, blob_name, local_file_path)
            chat_history = json.load(open(local_file_path, encoding="utf-8"))
            # Make sure the file is a dict
            if not isinstance(chat_history, dict):
                chat_history = {}
            # Optional: fix any inner structure if needed
            return chat_history
        return {}
    except Exception as e:
        st.error(f"Failed to load chat history: {str(e)}")
        return {}
    # try:
    #     # Check if the blob (file) exists in Azure Blob Storage.
    #     if file_exists_in_blob(blob_name):
    #         local_file_path = "chat_history.json"
    #         # Attempt to download the blob to a local file.
    #         if download_from_blob_storage(s3_bucket_name, blob_name, local_file_path):
    #             with open(local_file_path, encoding="utf-8") as file:
    #                 chat_history = json.load(file)
    #             # Ensure the loaded JSON is a dictionary.
    #             if not isinstance(chat_history, dict):
    #                 chat_history = {}
    #             return chat_history
    #     # If the file doesn't exist or download failed, return an empty dict.
    #     return {}
    # except Exception as e:
    #     st.error(f"Failed to load chat history: {str(e)}")
    #     return {}


def file_exists_in_blob(file_name):
    """Check if a file with the same name exists in S3."""
    try:
        s3_client.head_object(Bucket=s3_bucket_name, Key=file_name)
        return True
    except Exception as e:
        if e.response['Error']['Code'] == '404':
            return False
        else:
            raise e  # Re-raise other exceptions
    # """Check if a file with the same name exists in Azure Blob Storage."""
    # # Initialize BlobServiceClient
    # blob_service_client = BlobServiceClient.from_connection_string(connection_string)
    # container_client = blob_service_client.get_container_client(s3_bucket_name)

    # # Get the list of blobs in the container
    # blob_list = [blob.name for blob in container_client.list_blobs()]

    # # Check if the file name exists in the blob list
    # return file_name in blob_list

# Function to upload file to Azure Blob Storage
def upload_to_blob_storage(local_file_path, bucket_name, s3_key):
    try:
        with open(local_file_path, "rb") as data:
            s3_client.upload_fileobj(data, bucket_name, s3_key)

        # st.success(f"File '{s3_key}' successfully uploaded to S3.")
    except Exception as e:
        st.error(f"Failed to upload file to S3: {str(e)}")
    # try:
    #     blob_service_client = BlobServiceClient.from_connection_string(connection_string)
    #     blob_client = blob_service_client.get_blob_client(container=bucket_name, blob=s3_key)

    #     with open(local_file_path, "rb") as data:
    #         blob_client.upload_blob(data, overwrite=True)

    #     st.success(f"File '{s3_key}' successfully uploaded to Blob Storage.")
    # except Exception as e:
    #     st.error(f"Failed to upload file to Blob Storage: {str(e)}")

# Function to download file from Azure Blob Storage
def download_from_blob_storage(s3_bucket_name, s3_key, local_file_path):
    """Download a file from S3, or return False if not found."""
    try:
        with open(local_file_path, "wb") as file:
            s3_client.download_fileobj(s3_bucket_name, s3_key, file)
        return True
    except Exception as e:
        if e.response['Error']['Code'] == '404':
            print(f"File not found in S3: {s3_key}")
            return False
        else:
            print(f"Failed to download {s3_key}: {str(e)}")
            return False
    # try:
    #     blob_service_client = BlobServiceClient.from_connection_string(connection_string)
    #     blob_client = blob_service_client.get_blob_client(container=s3_bucket_name, blob=s3_key)

    #     with open(local_file_path, "wb") as file:
    #         file.write(blob_client.download_blob().readall())
    #     return True
    # except Exception as e:
    #     print(f"Failed to download {s3_key}: {str(e)}")
    #     return False

def create_word_doc(text):
    doc = Document()
    doc.add_heading("Chat Answer", level=1)
    doc.add_paragraph(text)
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Function to Generate titan embeddings
def generate_titan_embeddings(text):
    try:
        # Generate embeddings using MPNet
        embedding = mpnet_model.encode(text, normalize_embeddings=True)
        return np.array(embedding)
    except Exception as e:
        print(f"Error generating embeddings: {e}")
        return None  # Return None to handle errors gracefully

def call_llm_api(system_message, user_query):
    try:
        payload = {"system": system_message, "user": user_query, "max_tokens": 16384}
        r = requests.post(BACKEND_URL, json=payload, timeout=360)
        r.raise_for_status()
        return r.json()["answer"]
    
    except Exception as e:
        return f"An error occurred: {str(e)}"

# Faiss index initialization
dimension = 768  # Embedding dimension for text embeddings v3
faiss_index = faiss.IndexFlatL2(dimension)
metadata_store = []

# Updated `save_index_and_metadata` to upload files to Azure Blob Storage
def save_index_and_metadata():
    # Save files locally
    faiss.write_index(faiss_index, FAISS_INDEX_PATH)
    with open(METADATA_STORE_PATH, "wb") as f:
        pickle.dump(metadata_store, f)

    # Upload files to Blob Storage
    try:
        upload_to_blob_storage(FAISS_INDEX_PATH, s3_bucket_name, os.path.basename(FAISS_INDEX_PATH))
        upload_to_blob_storage(METADATA_STORE_PATH, s3_bucket_name, os.path.basename(METADATA_STORE_PATH))
    except Exception as e:
        print(f"Error uploading index or metadata to Blob Storage: {str(e)}")

# Function to load index and metadata
# Load index and metadata from Azure Blob Storage or initialize new
def load_index_and_metadata():
    global faiss_index, metadata_store

    index_blob_name = os.path.basename(FAISS_INDEX_PATH)
    metadata_blob_name = os.path.basename(METADATA_STORE_PATH)

    # Download files from Blob Storage if available
    index_downloaded = download_from_blob_storage(s3_bucket_name, index_blob_name, FAISS_INDEX_PATH)
    metadata_downloaded = download_from_blob_storage(s3_bucket_name, metadata_blob_name, METADATA_STORE_PATH)

    if index_downloaded and metadata_downloaded:
        # Load FAISS index and metadata store
        try:
            faiss_index = faiss.read_index(FAISS_INDEX_PATH)
            with open(METADATA_STORE_PATH, "rb") as f:
                metadata_store = pickle.load(f)
            print("Index and metadata loaded from Storage.")
        except Exception as e:
            print(f"Failed to load index or metadata: {str(e)}")
            # Initialize empty index and metadata if loading fails
            faiss_index = faiss.IndexFlatL2(dimension)
            metadata_store = []
    else:
        print("Index or metadata not found in Blob Storage. Initializing new.")
        # Initialize empty index and metadata
        faiss_index = faiss.IndexFlatL2(dimension)
        metadata_store = []

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    pages_text = []
    processing_message_placeholder = st.empty()
    progress_bar = st.progress(0)

    for page_num in range(len(doc)):
        processing_message_placeholder.write(f"Processing page {page_num + 1}/{len(doc)}...")
        temp_image_path = None

        try:
            # Render page as an image
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            temp_image_path = os.path.join(tempfile.gettempdir(), f"page_{page_num}.png")
            pix.save(temp_image_path)

            with open(temp_image_path, "rb") as image_file:
                image_bytes = image_file.read()

            # Use basic OCR (DetectDocumentText) instead of full analysis
            response = textract_client.detect_document_text(
                Document={'Bytes': image_bytes}
            )

            page_content = {"page_num": page_num + 1, "text": ""}
            
            # Extract text directly from OCR results
            text_lines = []
            for block in response.get('Blocks', []):
                if block['BlockType'] == 'LINE' and 'Text' in block:
                    text_lines.append(block['Text'])
            
            page_content["text"] = "\n".join(text_lines)
            pages_text.append((page_num + 1, page_content['text'], page_content['text']))

        except Exception as e:
            st.error(f"Error processing page {page_num + 1}: {str(e)}")
            continue

        finally:
            if temp_image_path and os.path.exists(temp_image_path):
                os.remove(temp_image_path)

        progress_bar.progress((page_num + 1) / len(doc))

    return pages_text

def chunk_text(pages_text, chunk_size=1):
    """Create non-overlapping chunks of text from pages."""
    chunks = []
    total_pages = len(pages_text)

    # Loop to create non-overlapping chunks
    for i in range(0, total_pages, chunk_size):
        chunk_parts = []
        for j in range(chunk_size):
            if i + j < total_pages:  # Ensure we do not exceed total pages
                page_num, text = pages_text[i + j]
                chunk_parts.append(f"The following text is from Page {page_num}:\n``````\n\n")
        
        if chunk_parts:  # Only append non-empty chunks
            chunks.append(''.join(chunk_parts))

    return chunks

def delete_file(file_name):
    try:
        # Delete from Amazon S3:
        s3_client.delete_object(Bucket=s3_bucket_name, Key=file_name)
        st.success(f"Deleted file '{file_name}' from S3 bucket.")
    except Exception as e:
        st.error(f"Error deleting file: {str(e)}")
    
    # Remove the file from metadata_store (filter out all records with that filename)
    global metadata_store
    metadata_store = [md for md in metadata_store if md["filename"] != file_name]
    
    # Optionally, update your index and re-save metadata
    save_index_and_metadata()

def user_has_file_access(username, file_name):
    """
    Checks whether the given username already has access to a file with file_name.
    This means the file is present in metadata_store with the same filename 
    AND the user is either the owner or is in that file's shared_with list.
    """
    for record in metadata_store:
        if record.get("filename") == file_name:
            if record.get("owner") == username or username in record.get("shared_with", []):
                return True
    return False

def add_pdf_to_index(pdf_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        temp_pdf.write(pdf_file.read())
        temp_pdf_path = temp_pdf.name
        
        pages_text = extract_text_from_pdf(temp_pdf_path)
        
        total_pages = len(pages_text)
        progress_bar = st.progress(0)  # Start at 0%
        
        # Create a placeholder for processing messages
        processing_message_placeholder = st.empty()

        # page_clf_text = ""
        
        for page_num, text, page_content in pages_text:
            processing_message_placeholder.write(f"Processing page {page_num}/{total_pages}...")

            embedding = generate_titan_embeddings(text)
            faiss_index.add(embedding.reshape(1, -1))
            metadata_store.append({
                "filename": os.path.basename(pdf_file.name),
                "page": page_num,
                "text": page_content,
                "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
                "shared_with": [] 
            })            
            
            # Update progress bar and clear previous message after updating
            progress_bar.progress(page_num / total_pages)  # Update to reflect current page
        
        save_index_and_metadata()
        
        try:
            # Upload the processed file to Azure Blob Storage
            blob_name = os.path.basename(pdf_file.name)
            upload_to_blob_storage(temp_pdf_path, s3_bucket_name, blob_name)
            
            os.remove(temp_pdf_path)  # Ensure file is removed after processing
        except PermissionError:
            st.warning("Could not delete the temporary PDF file. It might still be in use.")

def get_page_range_for_files(selected_files):
    """
    Finds the min and max page numbers for the selected files from metadata.
    """
    page_ranges = {}
    
    for file in selected_files:
        # Filter metadata for the selected file
        file_metadata = [metadata for metadata in metadata_store if metadata['filename'] == file]
        
        # Get the minimum and maximum page numbers for the file
        if file_metadata:
            min_page = min(metadata['page'] for metadata in file_metadata)
            max_page = max(metadata['page'] for metadata in file_metadata)
            page_ranges[file] = (min_page, max_page)
    
    return page_ranges

 
def query_documents_viz(selected_files, selected_page_ranges, query, top_k, web_search, llm_model):  
    query_embedding = generate_titan_embeddings(query).reshape(1, -1)
    if faiss_index.ntotal == 0:
        st.error("The FAISS index is empty. Please upload a PDF to populate the index.")
        return [], "No data available to query."

    # Fetch all metadata for the given query
    k = faiss_index.ntotal  # Initial broad search
    distances, indices = faiss_index.search(query_embedding, k)
    
    filtered_results = []
    for dist, idx in zip(distances[0], indices[0]):
        if idx < len(metadata_store):
            metadata = metadata_store[idx]
            if metadata['filename'] in selected_files:
                min_page, max_page = selected_page_ranges.get(metadata['filename'], (None, None))
                if min_page and max_page and min_page <= metadata['page'] <= max_page:
                    filtered_results.append((dist, idx))
    
    # Limit to topK after filtering
    top_k_results = sorted(filtered_results, key=lambda x: x[0])[:top_k]
    top_k_metadata = [metadata_store[idx] for _, idx in top_k_results]
    
    # Existing prompt modification
    # Existing prompt modification
    query_prompt = f"""  
    Given the extracted data from the uploaded documents, please respond to the user queries. 
    # Important: Remember these answers are for the leadership team so any valuable additional insights are always appreciated.
    # Always provide all necessary details and quote sources when getting yur answers.
    # Note : If the customer query requires a bar chart or graph, generate the equivalent Python code with all necessary imports 
    # and ensure the code uses Plotly for interactive graph creation (not Matplotlib).
    
    Extracted data from Index (In JSON Formatting): 
        ```
        {json.dumps(top_k_metadata)}
        ```
    """

    user_query = f"The User Question was: {query} \n\n"

    
    ws_response = ""

    if web_search:
        ws_query = query
        # Call the LLM API to get the answer
        # To install, run: pip install tavily-python


        client = TavilyClient(api_key=TAVILY_API)

        ws_response = client.search(
            query=ws_query,
            search_depth="advanced",
            include_raw_content=True
        )

        print(ws_response)

        wsp = f"""
        # Feel free to use the Web Search Results for Additional Context as well:

        {json.dumps(ws_response)}
        """
        if llm_model=="Skywork R1V2":
            answer = call_llm_api(query_prompt, user_query+wsp)
    else:
        if llm_model=="Skywork R1V2":
            answer = call_llm_api(query_prompt, user_query)

    return answer

###

def query_documents_with_page_range(
    selected_files: List[str],
    selected_page_ranges: Dict[str, Tuple[int, int]],
    prompt: str,
    top_k: int,
    last_messages: List[dict],
    web_search: bool,
    llm_model: str,
    draft_mode: bool,
    analyse_mode: bool,
    eco_mode: bool
):
    """
    An upgraded multi-step Deep Research workflow:
     1. Retrieve relevant doc chunks (FAISS).
     2. Perform iterative chunk-based analysis with repeated web searches (if analyse_mode=True).
     3. If draft_mode=True => produce a structured, chapter-wise "research paper" style doc.
        Else => produce a final Q&A style answer.
    
    Returns:
        (doc_chunks, final_text, ws_data)
        Where:
          - doc_chunks: the top-K relevant chunks from the documents
          - final_text: either a final "chapter-wise" doc or a Q&A answer
          - ws_data: final set of web search results (including any iterative expansions)
    """

    ##############################################################################
    # 1) Helper: Retrieve relevant doc chunks from FAISS
    ##############################################################################
    def retrieve_relevant_chunks(
        user_query: str,
        files: list,
        page_ranges: dict,
        topk: int
    ) -> list:
        """
        1) Embed the user query.
        2) Search FAISS for top-k results.
        3) Filter by selected files & page ranges.
        4) Return a list of chunk metadata (doc excerpt, filename, page, etc.).
        """
        if faiss_index.ntotal == 0:
            return []

        query_emb = generate_titan_embeddings(user_query)
        if query_emb is None:
            return []

        query_emb = query_emb.reshape(1, -1)
        total_docs = faiss_index.ntotal
        distances, indices = faiss_index.search(query_emb, total_docs)

        filtered = []
        for dist, idx in zip(distances[0], indices[0]):
            if idx < len(metadata_store):
                meta = metadata_store[idx]
                if meta["filename"] in files:
                    pg_min, pg_max = page_ranges.get(meta["filename"], (None, None))
                    if pg_min is not None and pg_max is not None:
                        if pg_min <= meta["page"] <= pg_max:
                            filtered.append((dist, idx))

        # Keep top_k best matches
        filtered_sorted = sorted(filtered, key=lambda x: x[0])[:topk]
        top_k_chunks = [metadata_store[x[1]] for x in filtered_sorted]
        return top_k_chunks

    ##############################################################################
    # 2) Helper: Web Search
    ##############################################################################
    def perform_web_search(query: str) -> dict:
        """
        Calls TavilyClient (or other web search) with the user query.
        Return raw web results as a dict.
        """
        from tavily import TavilyClient  # Ensure you have tavily installed
        client = TavilyClient(api_key=TAVILY_API)
        try:
            results = client.search(
                query=query,
                search_depth="advanced",
                include_raw_content=True
            )
            return results
        except Exception as e:
            return {"error": str(e), "query": query}

    ##############################################################################
    # 3) Iterative Analysis Agent
    ##############################################################################
    def analysis_agent_iterative(
        doc_chunks: list,
        initial_ws_data: dict,
        user_query: str,
        last_messages: list,
        eco_mode: bool,
        max_iterations: int = 5
    ) -> dict:
        """
        1) Summarize doc_chunks in groups.
        2) Merge them into a partial summary.
        3) Ask LLM if more web references are needed. If yes, do more queries.
        4) Keep refining until no more queries or max_iterations reached.
        5) Return final analysis text + references + all web data.

        The final text is 'analysis_summary'. 
        references: track chunk references
        all_web_data: store all web searches performed
        """

        # (A) Basic chunk-based summarization
        if not doc_chunks:
            return {
                "analysis_summary": "No relevant doc context found.",
                "references": [],
                "all_web_data": [],
                "web_search_queries": []
            }

        # Group chunks, e.g., 3 at a time
        chunk_size = 20
        grouped = [doc_chunks[i:i+chunk_size] for i in range(0, len(doc_chunks), chunk_size)]

        chunk_summaries = []
        references = []
        sys_msg = (
            "You are a deep-research analysis agent. For each chunk group, produce a thorough, "
            "multi-paragraph analysis referencing the filename/page. Then combine them all at the end."
        )

        for i, group in enumerate(grouped, start=1):
            group_text = "\n".join(
                f"File: {c['filename']}, Page: {c['page']} => {c['text']}"
                for c in group
            )
            user_prompt = f"CHUNK {i}:\n{group_text}\n\nPlease summarize in detail."

            result = call_llm_api(sys_msg, user_prompt)

            chunk_summaries.append(result)
            references.append({"section": i, "chunks": group})

        partial_summary = "\n\n".join(chunk_summaries)
        all_web_data = [initial_ws_data] if initial_ws_data else []
        web_search_queries = []

        # (B) Iteration loop to refine queries
        for iteration in range(max_iterations):
            # Step 1: Ask if more references needed
            refine_sys = (
                "You are an assistant that identifies missing references or clarifications "
                "from the partial analysis. If more web research is needed, propose up to 3 new search queries "
                "in a JSON array. If no more references are needed, return an empty array `[]`."
            )
            refine_user = f"""
            PARTIAL SUMMARY (Iteration {iteration+1}):
            {partial_summary}

            WEB SEARCH DATA SO FAR: {json.dumps(all_web_data, indent=2)}

            Return only JSON like:
            ```
            ["search query one", "search query two", ...]
            ```
            """

            refine_output = call_llm_api(refine_sys, refine_user)

            # Attempt to parse it as JSON
            try:
                # extract the JSON portion
                if "```" in refine_output:
                    inside = refine_output.split("```")[1].strip()
                    refined_queries = json.loads(inside)
                else:
                    refined_queries = json.loads(refine_output.strip())
            except:
                refined_queries = []

            if not isinstance(refined_queries, list) or len(refined_queries) == 0:
                # no additional queries
                break

            # Step 2: Perform new searches
            for q in refined_queries:
                wres = perform_web_search(q)
                web_search_queries.append(q)
                all_web_data.append({"query": q, "results": wres})

            # Step 3: Merge new web data into partial summary
            merge_sys = (
                "You are an analysis agent. Merge the original partial analysis below with new web data. "
                "Expand or refine it to be more comprehensive."
            )
            merge_user = f"""
            CURRENT ANALYSIS:
            {partial_summary}

            NEW WEB SEARCH RESULTS:
            {json.dumps(all_web_data[-len(refined_queries):], indent=2)}
            """

            new_summary = call_llm_api(merge_sys, merge_user)

            partial_summary = new_summary

        # Return final
        return {
            "analysis_summary": partial_summary,
            "references": references,
            "all_web_data": all_web_data,
            "web_search_queries": web_search_queries
        }

    ##############################################################################
    # 4) Chapter-Wise Drafting Agent (for `draft_mode=True`)
    ##############################################################################
    def drafting_agent_chapterwise(
        user_query: str,
        final_analysis_text: str,
        references: list,
        last_msgs: list,
        all_web_data: list,
        eco_mode: bool
    ) -> str:
        """
        Produces a "research paper" style structured doc:
         1) Table of Contents
         2) Chapters (Introduction, Main Analysis, etc.)
         3) Conclusion or Findings

        Returns a single large text/markdown.
        """
        sys_msg = (
            "You are a drafting agent. You have a final analysis text plus references. "
            "Produce a multi-chapter research paper with the following sections:\n\n"
            "1) Table of Contents\n"
            "2) Chapter 1: Introduction\n"
            "3) Chapter 2: Detailed Analysis\n"
            "4) Chapter 3: Findings/Observations\n"
            "5) Conclusion\n\n"
            "Each chapter should be carefully written in a professional legal research style. "
            "Include references to the doc pages or web data as footnotes or inline citations."
        )

        user_prompt = f"""
        USER PROMPT:
        {user_query}

        FINAL ANALYSIS TEXT:
        {final_analysis_text}

        REFERENCES:
        {json.dumps(references, indent=2)}

        WEB DATA:
        {json.dumps(all_web_data, indent=2)}

        LAST MESSAGES:
        {json.dumps(last_msgs, indent=2)}

        Please output a structured 'paper' with headings, subheadings, and a concluding section.
        """

        doc_text = call_llm_api(sys_msg, user_prompt)

        return doc_text

    ##############################################################################
    # 5) Q&A Agent (if draft_mode=False)
    ##############################################################################
    def qna_agent_simple(
        user_query: str,
        final_analysis_text: str,
        references: list,
        last_msgs: list,
        all_web_data: list,
        model_name: str
    ) -> str:
        """
        Return a simpler Q&A style answer referencing the final analysis and web data.
        """
        sys_msg = (
            "You are a Q&A agent. Use the final analysis text and references to answer succinctly. "
            "Provide references to doc pages or web data where relevant."
        )
        user_prompt = f"""
        USER PROMPT:
        {user_query}

        FINAL ANALYSIS TEXT:
        {final_analysis_text}

        REFERENCES:
        {json.dumps(references, indent=2)}

        WEB DATA:
        {json.dumps(all_web_data, indent=2)}

        LAST MESSAGES:
        {json.dumps(last_msgs, indent=2)}
        """

        return call_llm_api(sys_msg, user_prompt)

    ##############################################################################
    # MAIN LOGIC
    ##############################################################################

    # A) Retrieve doc chunks
    doc_chunks = retrieve_relevant_chunks(prompt, selected_files, selected_page_ranges, top_k)
    if not doc_chunks:
        return [], "No relevant documents found or FAISS index is empty.", {}

    # B) If analyse_mode => do iterative analysis with web searching
    #    else => minimal pass
    if analyse_mode:
        # Always do at least one initial web search with the raw prompt
        initial_ws_data = {}
        if web_search:
            initial_ws_data = perform_web_search(prompt)

        analysis_output = analysis_agent_iterative(
            doc_chunks=doc_chunks,
            initial_ws_data=initial_ws_data,
            user_query=prompt,
            last_messages=last_messages,
            eco_mode=eco_mode,
            max_iterations=5  # You can customize this
        )

        final_analysis_text = analysis_output["analysis_summary"]
        references = analysis_output["references"]
        all_web_data = analysis_output["all_web_data"]

    else:
        # minimal pass: combine doc chunks into one text
        final_analysis_text = "\n".join(ch["text"] for ch in doc_chunks)
        references = doc_chunks
        all_web_data = {}
        # If user specifically toggled web_search for non-analyse mode, we do a single search (optional)
        if web_search:
            all_web_data = perform_web_search(prompt)

    # C) Output
    if draft_mode:
        # Produce a "research paper" style doc
        doc_text = drafting_agent_chapterwise(
            user_query=prompt,
            final_analysis_text=final_analysis_text,
            references=references,
            last_msgs=last_messages,
            all_web_data=all_web_data,
            eco_mode=eco_mode
        )
        return doc_chunks, doc_text, all_web_data
    else:
        # Simple Q&A
        answer = qna_agent_simple(
            user_query=prompt,
            final_analysis_text=final_analysis_text,
            references=references,
            last_msgs=last_messages,
            all_web_data=all_web_data,
            model_name=llm_model
        )
        return doc_chunks, answer, all_web_data

###

def final_format(top_k_metadata, answer, ws_response):
    sys_msg = """
    You are a helpful Legal Assistant. 
    You Specialise in Formatting Generated Answers.

    """
    
    input_context = f"""
    Below is the Generated Answer Fetched from LLM:
    <<<{answer}>>>

    The top K most relevant contexts fetched from the documents are as follows:
        {json.dumps(top_k_metadata, indent=4)}

    ##########################################################################
        
    The Web Search Results Fetched are as follows:
        {json.dumps(ws_response)}

    ##########################################################################
        
    """

    final_op_format = '''
        For Transparency and Explanability, we Would like you to do the following:
            1. Break the Answer into Logical Shards.
            2. For Each Shard, Map them to the Relevant Sources (From the Provided TopK Context) in the Formatting as Presented below
        
        Note: The Shards when concatenated should help us regenerate the provided `Generated Answer`
        Approach the Above Step by Step.

        # Final Output Format:
        ```
        {
            "segmented_answer":    
                [
                    {
                        "section": "The first shard of the generated answer",
                        "sources": [
                                        {"filename": "The First Filename", "page": "page_num", "text": "The Relevant text from the page only"},
                                        {"filename": "The Second Filename", "page": "page_num", "text": "The Relevant text from the page only"},
                                        . # Add more sources if necessary
                                    ]
                    },
                    {
                        "section": "The second part of the generated answer", 
                        "sources": [
                                        {"filename": "The First Filename", "page": "page_num", "text": "The Relevant text from the page only"},
                                        {"filename": "The Second Filename", "page": "page_num", "text": "The Relevant text from the page only"},
                                        . # Add more sources if necessary
                                    ]
                    },
                    .
                    .
                ]
        }
        ```
    '''
    # Call the LLM API to get the answer
    answer = call_llm_api(sys_msg, input_context+final_op_format)
    try:
        # return json.loads(answer[7:-3])
        return json.loads(answer.split("```json")[1].split("```")[0])
    except:
        # return json.loads(answer[3:-3])
        return json.loads(answer.split("```")[1].split("```")[0])
    
def summarize_document_pages(filename, start_page, end_page, summary_prompt, eco_mode=False):
    """
    Summarize document pages for a given file using overlapping chunks if needed.
    Retrieves text from the metadata_store for pages in [start_page, end_page].
    """
    # Retrieve texts for the specified file and pages.
    pages = [metadata for metadata in metadata_store 
             if metadata['filename'] == filename and start_page <= metadata['page'] <= end_page]
    # Sort pages by page number.
    pages = sorted(pages, key=lambda x: x['page'])
    total_pages = len(pages)
    full_text = json.dumps(pages) # "\n".join(page["text"] for page in pages)

    # If the selected pages are less than 20, summarize in one go.
    if total_pages < 50:
        user_query = f"Summarize the following document text from Document {filename} \n{full_text}"
        summary = call_llm_api(summary_prompt, user_query)
        return summary

    # Else, create overlapping chunks.
    # Define base chunk size and overlap (e.g., 20 pages with 25% overlap => 5 pages)
    base_chunk_size = 50
    overlap = int(base_chunk_size * 0.25)  # 5 pages

    # Build chunks using a sliding window approach.
    chunk_summaries = []
    i = 0
    while i < total_pages:
        # Determine the chunk's pages (ensure we don't exceed total_pages)
        start_idx = i
        end_idx = min(i + base_chunk_size, total_pages)
        # chunk_text = json.dumps(pages[j] for j in range(start_idx, end_idx))
        chunk_text = json.dumps([pages[j] for j in range(start_idx, end_idx)])
        user_query = f"Summarize the following document from file `{filename}` text:\n{chunk_text}"
        chunk_summary = call_llm_api(summary_prompt, user_query)
        chunk_summaries.append(chunk_summary)
        # Advance by base_chunk_size minus the overlap.
        i += (base_chunk_size - overlap)
    
    # Consolidate the chunk summaries into a final summary.
    consolidation_prompt = summary_prompt + "\n\nPlease consolidate the following summaries into one overall summary:"
    consolidation_input = json.dumps(chunk_summaries)
    final_summary = call_llm_api(consolidation_prompt, consolidation_input)
    return final_summary

def get_web_recommendations(document_summaries, insights):
        # print(selected_files)
    # print(selected_page_ranges)
    qp_prompt = {
        "system_message": "You are an intelligent query refiner. Your job is to take the Document Summaries and Key observations (which may contain poor grammar or informal language) and generate a well-formed prompt for web search. The web search prompt should refine the query further to fetch relevant legal resources/ verdicts/ cases/ decisions online. Output only a JSON object with 'web_search_prompt' as key.",
        "user_query": f"Document Summaries: {json.dumps(document_summaries)}\n\nInsights: {insights}\n\nGenerate the JSON output with the improved prompt."
    }

    op_format = '''
    You are an expert legal researcher. Based on the above document summary and key observations, please craft a concise web search query that will help locate similar legal cases and corresponding verdicts. Your output should include:

        Specific legal terms or phrases that capture the core issues.
        Relevant jurisdiction or court references, if applicable.
        Any additional keywords that may refine the search (e.g., precedent case names or statutory citations).

    # Output Format:
    
    ```json
    {
        "web_search_prompt": "Further refined query designed to fetch relevant legal resources/ verdicts/ cases/ decisions from the web."
    }
    ```
    '''

    prompts = call_llm_api(qp_prompt["system_message"], qp_prompt["user_query"]+op_format)
    print(prompts)
    try:
        # return json.loads(answer[7:-3])
        prompt_op = json.loads(prompts.split("```json")[1].split("```")[0])
    except:
        # return json.loads(answer[3:-3])
        try:
            prompt_op = json.loads(prompts)
        except:
            prompt_op = json.loads(prompts.split("```")[1].split("```")[0])
    print(prompt_op)
    ws_query = prompt_op["web_search_prompt"]
    # Call the LLM API to get the answer
    # To install, run: pip install tavily-python


    client = TavilyClient(api_key=TAVILY_API)

    ws_response = client.search(
        query=ws_query,
        search_depth="advanced",
        include_raw_content=True
    )

    print(ws_response)

    web_response = call_llm_api(ws_prompt_lib, f"The Document Summaries: {json.dumps(document_summaries)} \n\n The Key Observations: {insights} \n\n The Web Search Results: {json.dumps(ws_response)}")
    nexus_insights = call_llm_api(insights_prompt_lib, f"The Document Summaries: {json.dumps(document_summaries)} \n\n The Key Observations: {insights} \n\n The Web Search Results: {json.dumps(ws_response)}")
    faq = call_llm_api(qna_prompt, f"The Document Summaries: {json.dumps(document_summaries)} \n\n The Key Observations: {insights} \n\n The Web Search Results: {json.dumps(ws_response)}")
    
    return web_response, nexus_insights, faq

# Function to generate post text (placeholder; replace with your actual API call)
def generate_post_text(task_type, template_instructions, description):
    prompt = f"Task Type: {task_type}\nTemplate: {template_instructions}\nDescription: {description}\nGenerate a creative post:"
    user_instructions = "Use your own Creative Liberties. Follow the instructions above and DO not Insert Text on the Images. Include Emojis and Stickers if Appropriate"

    # Dummy response for now:
    return call_llm_api(prompt, user_instructions)

    image_url = json.loads(result.model_dump_json())['data'][0]['url']
    return image_url

def sanitize_text(message_text):
    # Remove newline and tab characters
    sanitized = re.sub(r'[\n\t]', ' ', message_text)
    # Replace four or more consecutive spaces with a single space
    sanitized = re.sub(r' {4,}', ' ', sanitized)
    return sanitized.strip()

# WhatsApp sharing function (placeholder; replace with your actual API implementation)
def send_whatsapp_message(recipient, message_text):
    message_text = sanitize_text(message_text)

    url = "https://graph.facebook.com/v20.0/472581919270525/messages"
    headers = {
        "Content-Type": "application/json",
        "Authorization": WHATSAPP_TOKEN
    }
    data = {
        "messaging_product": "whatsapp",
        "recipient_type": "individual",
        "to": recipient,
        "type": "template",
        "template": {
            "name": "mytesttemplate",
            "language": {"code": "en_us"},
            "components": [
                {"type": "body", "parameters": [{"type": "text", "text": message_text}]}
            ]
        }
    }
    response = requests.post(url, headers=headers, json=data)
    return response.json()

# Email sending function from email.py (adjust if needed)
def send_email(subject, body, recipient):
    try:
        EMAIL_ADDRESS = EMAIL_ID  # Replace with your sender email
        EMAIL_PASSWORD = EMAIL_PWD  # Replace with your app password

        msg = MIMEMultipart()
        msg['From'] = EMAIL_ADDRESS
        msg['To'] = recipient
        msg['Subject'] = subject
        msg.attach(MIMEText(body, 'plain'))

        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
            server.send_message(msg)
        return {"status": "sent"}
    except Exception as e:
        return {"status": "error", "error": str(e)}

# --- Helper functions to load contacts ---

def load_whatsapp_contacts():
    # Replace with code to load from a JSON file if needed
    return {"Anubhav": "919874454959", "Prithviraj": "917980757702"}

def load_email_contacts():
    # Replace with code to load from a JSON file if needed
    return {"Anubhav": "anubhav@nekko.tech", "Prithviraj": "prithvi@nekko.tech"}

# Function to handle user input with text_area and predefined prompts
def user_input():

    prompt_options = list(prompt_library.keys())
    selected_prompt = st.selectbox("Select a predefined prompt:", prompt_options, index=0)
    
    # Auto-fill text area if a prompt is selected
    default_text = prompt_library[selected_prompt]
    user_message = st.text_area("Enter your message:", value=default_text, height=150)
    ol1, ol2, ol3, ol4, ol5, ol6, ol7, ol8 = st.columns(8)
    with ol8:
    # Submit button
       submit = st.button("Send")
    if submit and user_message.strip():
        default_text = "custom"
        st.session_state.user_message = ""  # Clear the text input
        return user_message.strip()
    return None

USERS = load_dict_from_json(users_file)
credentials = {"usernames": {}}
for user in USERS:
    # Assume each user has keys: 'username', 'name', and 'password'
    credentials["usernames"][user] = {
        "name": user,
        "password": USERS[user]
    }

def login():
    display_logo()
    st.title("Ready to Dive In? Sign In!")
    username = st.text_input("Username")
    password = st.text_input("Password", type="password")
    
    if st.button("Login"):
        if username in USERS and password == USERS[username]:
            # Mark the user as authenticated
            st.session_state["authenticated"] = True
            st.session_state["username"] = username
            st.session_state["logged_in"] = True

            # Set a persistent cookie
            cookies["username"] = username
            cookies.save()
            
            # Create the authenticator instance and store it in session state.
            # (Ensure you have a 'credentials' variable ready.)
            Authenticator = stauth.Authenticate(credentials, cookie_name='nexusdms/', key='abcdefgh', cookie_expiry_days=0)
            st.session_state["Authenticator"] = Authenticator
            
            st.success("Login successful!")
            # Optionally, you can call st.rerun() here if needed.
        else:
            st.error("Invalid username or password.")

def logout():
    # Retrieve the stored authenticator instance
    authenticator = st.session_state.get("Authenticator")
    if authenticator is not None:
        try:
            # Attempt logout (this might raise a KeyError if cookie already removed)
            logout_button = authenticator.logout('Log Out', 'sidebar')
        except KeyError:
            logout_button = True  # If cookie already removed, treat as successful logout.
        except Exception as err:
            st.error(f'Unexpected exception during logout: {err}')
            return
    else:
        logout_button = True

    if logout_button:
        # Update session state to reflect logout
        st.session_state["logged_out"] = True
        st.session_state["logged_in"] = False
        st.session_state["authenticated"] = False
        st.session_state["username"] = ""
        st.session_state["Authenticator"] = None
        
        # Clear the cookie as well.
        if "username" in cookies:
            # del cookies["username"]
            cookies["username"] = ""
            cookies.save()
        st.rerun()

def process_pdf(pdf_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
        temp_pdf.write(pdf_file.read())
        temp_pdf_path = temp_pdf.name
        
    pages_text = extract_text_from_pdf(temp_pdf_path)
    total_pages = len(pages_text)
    progress_bar = st.progress(0)
    processing_message_placeholder = st.empty()
    
    for page_num, text, page_content in pages_text:
        processing_message_placeholder.write(f"Processing page {page_num}/{total_pages}...")
        embedding = generate_titan_embeddings(text)
        faiss_index.add(embedding.reshape(1, -1))
        metadata_store.append({
            "filename": os.path.basename(pdf_file.name),
            "page": page_num,
            "text": page_content,
            "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
            "shared_with": [] 
        })
        progress_bar.progress(page_num / total_pages)
    
    save_index_and_metadata()
    try:
        blob_name = os.path.basename(pdf_file.name)
        upload_to_blob_storage(temp_pdf_path, s3_bucket_name, blob_name)
        os.remove(temp_pdf_path)
    except PermissionError:
        st.warning("Could not delete the temporary PDF file. It might still be in use.")

def extract_text_from_image(file_path):
    # Open image with fitz (PyMuPDF supports image reading too)
    img_doc = fitz.open(file_path)
    text_lines = []
    try:
        page = img_doc[0]  # Assume the entire image is one "page"
        pix = page.get_pixmap()
        temp_image_path = os.path.join(tempfile.gettempdir(), "temp_image.png")
        pix.save(temp_image_path)
        with open(temp_image_path, "rb") as image_file:
            image_bytes = image_file.read()
        response = textract_client.detect_document_text(Document={'Bytes': image_bytes})
        for block in response.get('Blocks', []):
            if block['BlockType'] == 'LINE' and 'Text' in block:
                text_lines.append(block['Text'])
        os.remove(temp_image_path)
    except Exception as e:
        st.error(f"Error processing image: {str(e)}")
    return "\n".join(text_lines)

def process_image(image_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_img:
        temp_img.write(image_file.read())
        temp_img_path = temp_img.name
    text = extract_text_from_image(temp_img_path)
    embedding = generate_titan_embeddings(text)
    faiss_index.add(embedding.reshape(1, -1))
    metadata_store.append({
        "filename": os.path.basename(image_file.name),
        "page": 1,  # Only one "page" for an image
        "text": text,
        "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
        "shared_with": [] 
    })
    save_index_and_metadata()
    try:
        blob_name = os.path.basename(image_file.name)
        upload_to_blob_storage(temp_img_path, s3_bucket_name, blob_name)
        os.remove(temp_img_path)
    except PermissionError:
        st.warning("Could not delete the temporary image file. It might still be in use.")

def process_docx(docx_file):
    # Save file locally
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_docx:
        temp_docx.write(docx_file.read())
        temp_docx_path = temp_docx.name

    try:
        doc = Document(temp_docx_path)
        # If page information is not available, you can split the document into chunks by paragraphs.
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip() != ""])
        # Optional: use a heuristic to approximate pages (e.g., split full_text into equal parts)
        chunks = [full_text[i:i+1000] for i in range(0, len(full_text), 1000)]
        for idx, chunk in enumerate(chunks, 1):
            embedding = generate_titan_embeddings(chunk)
            faiss_index.add(embedding.reshape(1, -1))
            metadata_store.append({
                "filename": os.path.basename(docx_file.name),
                "page": idx,
                "text": chunk,
                "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
                "shared_with": [] 
            })
    except Exception as e:
        st.error(f"Error processing DOCX: {str(e)}")
    finally:
        os.remove(temp_docx_path)
    save_index_and_metadata()

def process_pptx(pptx_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
        temp_pptx.write(pptx_file.read())
        temp_pptx_path = temp_pptx.name
    try:
        prs = Presentation(temp_pptx_path)
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            embedding = generate_titan_embeddings(slide_text)
            faiss_index.add(embedding.reshape(1, -1))
            metadata_store.append({
                "filename": os.path.basename(pptx_file.name),
                "page": idx,  # Slide number
                "text": slide_text,
                "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
                "shared_with": [] 
            })
    except Exception as e:
        st.error(f"Error processing PPTX: {str(e)}")
    finally:
        os.remove(temp_pptx_path)
    save_index_and_metadata()

def process_spreadsheet(file_obj):
    ext = os.path.splitext(file_obj.name)[1].lower()
    try:
        if ext == ".csv":
            df = pd.read_csv(file_obj)
            sheet_name = "csv"
        else:
            xls = pd.ExcelFile(file_obj)
            # Process all sheets or just the first one; here we use the first sheet
            sheet_name = xls.sheet_names[0]
            df = pd.read_excel(xls, sheet_name=sheet_name)
    except Exception as e:
        st.error(f"Error reading spreadsheet: {str(e)}")
        return

    chunk_size = 50
    num_chunks = (len(df) // chunk_size) + int(len(df) % chunk_size != 0)
    for i in range(num_chunks):
        chunk_df = df.iloc[i*chunk_size : (i+1)*chunk_size]
        chunk_text = chunk_df.to_string(index=False)
        embedding = generate_titan_embeddings(chunk_text)
        faiss_index.add(embedding.reshape(1, -1))
        metadata_store.append({
            "filename": f"{os.path.basename(file_obj.name)} ({sheet_name})",
            "page": i + 1,  # Store as an integer instead of "chunk_{i+1}"
            "text": chunk_text,
            "owner": st.session_state.get("username", "unknown"),  # Save the uploader’s username
            "shared_with": [] 
        })
    save_index_and_metadata()

def add_file_to_index(uploaded_file):
    ext = os.path.splitext(uploaded_file.name)[1].lower()
    
    if ext == ".pdf":
        process_pdf(uploaded_file)
    elif ext in [".jpg", ".jpeg", ".png"]:
        process_image(uploaded_file)
    elif ext in [".doc", ".docx"]:
        process_docx(uploaded_file)
    elif ext == ".pptx":
        process_pptx(uploaded_file)
    elif ext in [".xlsx", ".csv"]:
        process_spreadsheet(uploaded_file)
    else:
        st.error(f"Unsupported file type: {ext}")

# --- Helper Function: Overlay Logo and Dynamic Text ---
def load_image_from_input(image_input):
    """
    Loads an image from a local file path or a URL.
    """
    if isinstance(image_input, str) and image_input.startswith("http"):
        try:
            response = requests.get(image_input)
            response.raise_for_status()  # Ensure we got a valid response
            return Image.open(BytesIO(response.content))
        except requests.RequestException as e:
            st.error(f"Error loading image from URL: {e}")
            return None
    else:
        return Image.open(image_input)

def overlay_logo_and_text(image_input, logo_path, bottom_text, overlay_text, overlay_x_pct, overlay_y_pct, text_mode, font_adjuster):
    # Load main image from URL or local file
    main_image = load_image_from_input(image_input)
    if main_image is None:
        return None  # Stop processing if image load failed

    # Load logo (assuming logo is a local file)
    logo = Image.open(logo_path)

    # Get image dimensions
    img_width, img_height = main_image.size

    # Resize logo relative to image size (10% of image width)
    logo_width = int(img_width * 0.1)
    logo_height = int(logo.size[1] * (logo_width / logo.size[0]))  # Maintain aspect ratio
    logo = logo.resize((logo_width, logo_height), Image.Resampling.LANCZOS)

    # Paste the logo at top-left with 5% padding
    logo_x = int(img_width * 0.05)
    logo_y = int(img_height * 0.05)
    main_image.paste(logo, (logo_x, logo_y), logo)

    # Convert image to OpenCV format for text overlay
    image_cv = np.array(main_image)
    image_cv = cv2.cvtColor(image_cv, cv2.COLOR_RGB2BGR)

    # Define font settings (scaling font and thickness with image width)
    font = cv2.FONT_HERSHEY_SIMPLEX
    font_scale = max(1, img_width // 500) * font_adjuster
    thickness = max(2, img_width // 400)

    # Set text color based on user selection
    if text_mode.lower() == "light":
        font_color = (255, 255, 255)  # White text for light backgrounds
    else:
        font_color = (0, 0, 0)        # Black text for dark backgrounds

    # Add bottom text at 5% padding from left and 5% above the bottom edge
    bottom_x = int(img_width * 0.05)
    bottom_y = img_height - int(img_height * 0.05)
    cv2.putText(image_cv, bottom_text, (bottom_x, bottom_y), font, font_scale, font_color, thickness, cv2.LINE_AA)

    # Calculate overlay text position using percentage values (0 to 100)
    overlay_x = int(img_width * overlay_x_pct / 100)
    overlay_y = int(img_height * overlay_y_pct / 100)
    cv2.putText(image_cv, overlay_text, (overlay_x, overlay_y), font, font_scale, font_color, thickness, cv2.LINE_AA)

    # Convert back to PIL format
    final_image = Image.fromarray(cv2.cvtColor(image_cv, cv2.COLOR_BGR2RGB))
    return final_image

def main():
    # Try to restore authentication from session state or cookie.
    if not st.session_state["authenticated"]:
        cookie_username = cookies.get("username")
        if cookie_username != "" and cookie_username is not None:
            st.session_state["authenticated"] = True
            st.session_state["username"] = cookie_username
            st.session_state["logged_in"] = True

    if not st.session_state["authenticated"]:
        login()
        return

    display_logo()
    st.title("Document Query Assistant")

    # Initialize session state variables if not present.
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'current_conversation' not in st.session_state:
        st.session_state.current_conversation = None
    if 'sources' not in st.session_state:
        st.session_state.sources = []
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = load_chat_history()
    # Initialize file selection state.
    if 'selected_files' not in st.session_state:
        st.session_state.selected_files = []
    if 'selected_page_ranges' not in st.session_state:
        st.session_state.selected_page_ranges = {}
    if 'file_summaries' not in st.session_state:
        st.session_state.file_summaries = {}
    if "rename_mode" not in st.session_state:
        st.session_state["rename_mode"] = None
    if "share_chat_mode" not in st.session_state:
        st.session_state["share_chat_mode"] = False
    if "share_chat_conv" not in st.session_state:
        st.session_state["share_chat_conv"] = None
    if "share_chat_conv_id" not in st.session_state:
        st.session_state["share_chat_conv_id"] = None

    available_usernames = list(USERS.keys())
    
    current_user = st.session_state["username"]
    load_index_and_metadata()
    st.sidebar.header(f"Hello `{current_user}`")
    if st.sidebar.button("Log Out"):
        logout()  # Display the logout button in the sidebar

    st.sidebar.header("Options")
    option = st.sidebar.selectbox("Choose an option", ["Query Documents", "Query Advanced", "Upload Documents", "File Manager", "Usage Monitoring"])

    if option == "Upload Documents":
        st.header("Upload Documents")
        uploaded_files = st.file_uploader(
            "Upload one or more documents",
            type=["pdf", "jpg", "jpeg", "png", "doc", "docx", "pptx", "xlsx", "csv"],
            accept_multiple_files=True
        )

        if uploaded_files:
            for uploaded_file in uploaded_files:
                st.write(f"Processing {uploaded_file.name}...")
                current_user = st.session_state.get("username", "unknown")
                # if file_exists_in_blob(uploaded_file.name):
                #     st.warning(f"File '{uploaded_file.name}' already exists in Storage. Skipping upload.")
                # else:
                #     add_file_to_index(uploaded_file)
                #     st.success(f"File '{uploaded_file.name}' has been successfully uploaded and added to the index.")
                # Check if current user already has access
                if user_has_file_access(current_user, uploaded_file.name):
                    st.warning(f"File '{uploaded_file.name}' is already in your library (owned or shared). Skipping upload.")
                else:
                    # Proceed with normal upload & indexing
                    add_file_to_index(uploaded_file)
                    st.success(f"File '{uploaded_file.name}' has been successfully uploaded and added to your library.")

    elif option == "File Manager":
        # available_usernames = list(USERS.keys())
        st.header("My Uploaded Files")
        # Filter files for the current user (see user isolation in step 2)
        current_user = st.session_state.get("username", "unknown")
        available_files = list({
            md["filename"] 
            for md in metadata_store 
            if md.get("owner") == current_user or current_user in md.get("shared_with", [])
        })
        if available_files:
            if available_files:
                for i, fname in enumerate(available_files):
                    col1, col2 = st.columns([0.7, 0.3])
                    with col1:
                        st.write(fname)
                    with col2:
                        if st.button("Delete", key=f"del_{fname}_{i}"):
                            delete_file(fname)
        else:
            st.sidebar.info("No files uploaded yet.")

        # In the File Manager section:
        st.sidebar.header("Share a File")
        file_to_share = st.sidebar.selectbox("Select a file to share", available_files)
        share_with = st.sidebar.multiselect("Select user(s) to share with", options=available_usernames)

        if st.sidebar.button("Share File"):
            # Update metadata for the file if the current user is the owner
            for md in metadata_store:
                if md["filename"] == file_to_share and md.get("owner") == current_user:
                    md.setdefault("shared_with", []).extend(share_with)
                    md["shared_with"] = list(set(md["shared_with"]))  # Remove duplicates
                    st.success(f"Shared {file_to_share} with {', '.join(share_with)}")
            save_index_and_metadata()

    elif option == "Query Documents":
        st.header("Query Documents")
        st.sidebar.header("Settings")
        llm_model = st.sidebar.selectbox("Choose Your Model", ["Skywork R1V2"])

        # "New Chat" button resets conversation and state.
        if st.sidebar.button("New Chat"):
            st.session_state.current_conversation_id = None  # Clear any old conversation ID
            st.session_state.current_conversation = None
            st.session_state.messages = []
            st.session_state.sources = []
            # Clear previous file/page selections too.
            st.session_state.selected_files = []
            st.session_state.selected_page_ranges = {}
            st.success("Started a new conversation.")
        eco_mode = st.sidebar.toggle("Eco Mode", value=False)
        web_search = st.sidebar.toggle("Enable Web Search")
        draft_mode = st.sidebar.toggle("Enable Draft Mode (To Generate Documents/ Arguments)")
        analyse_mode = st.sidebar.toggle("Enable Analyse Mode (For Deeper Analysis and Search) [Consumes more Tokens]")
        top_k = st.sidebar.slider("Select Top-K Results", min_value=1, max_value=100, value=50, step=1)

        # File and Page Range Selection
        # available_files = list(set([metadata['filename'] for metadata in metadata_store]))
        current_user = st.session_state.get("username", "unknown")
        # Only include files where the owner is the current user or shared with the user.
        available_files = list({
            md["filename"] 
            for md in metadata_store 
            if md.get("owner") == current_user or current_user in md.get("shared_with", [])
        })


        if available_files:
            # Use multiselect and store the selection in session state.
            st.session_state.selected_files = st.multiselect(
                "Select files to include in the query:",
                available_files,
                default=st.session_state.selected_files
            )
            if len(st.session_state.selected_files) > 4:
                st.warning("For best results, select a maximum of 4 files.")
                # return

            page_ranges = get_page_range_for_files(st.session_state.selected_files)
            selected_page_ranges = {}
            # For each file, show page range inputs and store values in session state.
            for file in st.session_state.selected_files:
                min_page, max_page = page_ranges[file]
                col1, col2 = st.sidebar.columns(2)
                with col1:
                    start_page = st.number_input(
                        f"Start page for {file}",
                        min_value=min_page,
                        max_value=max_page,
                        value=page_ranges[file][0],
                        key=f"start_{file}"
                    )
                with col2:
                    end_page = st.number_input(
                        f"End page for {file}",
                        min_value=min_page,
                        max_value=max_page,
                        value=page_ranges[file][1],
                        key=f"end_{file}"
                    )
                selected_page_ranges[file] = (start_page, end_page)
            st.session_state.selected_page_ranges = selected_page_ranges


        summaries = {}
        # selected_files = []
        # At the top of the "Query Advanced" option, display document summaries.
        if st.session_state.selected_files and st.button("Nekko Insignts"):
            st.subheader("Document Summaries")
            # Create columns (one per selected file)
            # cols = st.columns(len(selected_files))
            for idx, file in enumerate(st.session_state.selected_files):
                # Get the page range for the file.
                min_page, max_page = selected_page_ranges.get(file, (None, None))
                if min_page is not None and max_page is not None:
                    # Use the summary prompt as the system message (you may have a variable 'summary_prompt' already defined).
                    summary = summarize_document_pages(file, min_page, max_page, summary_prompt)
                    summaries[file] = summary
                    # with cols[idx]:
                    st.markdown(f"**{file} (Pages {min_page}-{max_page}) Summary:**")
                    with st.expander("Click to view"):
                        st.write(summary)
            
            # cols1, cols2 = st.columns(2)
            st.subheader("Observations and Insights")
            # with cols1:
            st.markdown("**Key Observations**")
            insights = call_llm_api(insights_prompt, json.dumps(summaries))
            with st.expander("Click to view"):
                st.write(insights)
            # with cols2:
            st.markdown("**Web Search Results**")
            web_response, nexus_insights, faq = get_web_recommendations(summaries, insights)
            with st.expander("Click to view"):
                st.write(web_response)

            # cols3, cols4 = st.columns(2)
            st.subheader("Nexus Intelligence")
            # with cols3:
            st.markdown("**Nexus Intelligent Insights**")
            with st.expander("Click to view"):
                st.write(nexus_insights)
            # with cols4:
            st.markdown("**Frequently Asked Questions**")
            with st.expander("Click to view"):
                st.write(faq)

        st.sidebar.header("Previous Conversations")
        user_conversations = st.session_state.chat_history.get(current_user, [])

        # unique_conversations = []
        # seen_labels = {}
        # for conv in user_conversations:
        #     conv_label = conv.get("label") or conv.get('messages', [{}])[0].get("content", "")[:50]
        #     if conv_label not in seen_labels:
        #         seen_labels[conv_label] = conv
        #         unique_conversations.append(conv)
        # unique_conversations.sort(key=lambda x: x.get("timestamp", ""), reverse=True)

        # NEW CODE: Sort the entire list by timestamp
        unique_conversations = sorted(
            user_conversations, 
            key=lambda x: x.get("timestamp", ""), 
            reverse=True
        )

        for conv in unique_conversations:
            # Use the conversation's timestamp as a unique identifier.
            conv_id = conv.get("timestamp")
            default_label = conv.get("label") or conv.get('messages', [{}])[0].get("content", "")[:50]
            
            # Check if this conversation is in rename mode.
            if st.session_state.get("rename_mode") == conv_id:
                new_label = st.sidebar.text_input("Rename Conversation", value=default_label, key=f"rename_input_{conv_id}")
                if st.sidebar.button("Save", key=f"save_rename_{conv_id}"):
                    user = st.session_state.username
                    if user in st.session_state.chat_history:
                        for idx, stored_conv in enumerate(st.session_state.chat_history[user]):
                            if stored_conv.get("timestamp") == conv_id:
                                # Only update the label
                                st.session_state.chat_history[user][idx]["label"] = new_label
                                break 
                    save_chat_history(st.session_state.chat_history)
                    st.session_state["rename_mode"] = None
                    st.sidebar.success("✅ Conversation renamed successfully!")
                    st.rerun()
            else:
                col1, col2, col3, col4 = st.sidebar.columns([0.5, 0.2, 0.2, 0.1])
                if col1.button(default_label, key=f"load_{conv_id}"):
                    st.session_state.current_conversation = conv
                    st.session_state.messages = conv.get('messages', [])
                    st.session_state.selected_files = conv.get('files', [])
                    st.session_state.selected_page_ranges = conv.get('page_ranges', {})
                    st.session_state.current_conversation_id = conv_id
                    st.rerun()
                if col2.button("✏️", key=f"rename_button_{conv_id}"):
                    st.session_state["rename_mode"] = conv_id
                    st.rerun()
                if col3.button("🗑️", key=f"delete_{conv_id}"):
                    st.session_state["confirm_delete_conv"] = conv
                    st.rerun()
                if col4.button("📤", key=f"share_chat_{conv_id}"):
                    st.session_state["share_chat_conv"] = conv
                    st.session_state["share_chat_conv_id"] = conv_id
                    st.session_state["share_chat_mode"] = True
                    st.rerun()

        # Insert the share-chat snippet below the conversation list:
        if st.session_state.get("share_chat_mode"):
            st.header("Share Chat Conversation")
            share_chat_with = st.multiselect("Select user(s) to share with", options=available_usernames)
            if st.button("Confirm Share Chat"):
                chat_to_share = st.session_state["share_chat_conv"]
                # For each target user, append a deep copy of the conversation
                for user in share_chat_with:
                    if user in st.session_state.chat_history:
                        st.session_state.chat_history[user].append(copy.deepcopy(chat_to_share))
                    else:
                        st.session_state.chat_history[user] = [copy.deepcopy(chat_to_share)]
                save_chat_history(st.session_state.chat_history)
                st.success("Chat conversation shared successfully!")
                # Reset share mode variables
                st.session_state["share_chat_mode"] = False
                st.session_state.pop("share_chat_conv", None)
                st.session_state.pop("share_chat_conv_id", None)
                st.rerun()

        # If a conversation is marked for deletion, confirm deletion.
        if "confirm_delete_conv" in st.session_state:
            chat_name = (
                st.session_state["confirm_delete_conv"].get("label")
                or st.session_state["confirm_delete_conv"].get('messages', [{}])[0].get("content", "")[:50]
            )
            st.warning(f"Are you sure you want to delete '{chat_name}' conversation?")
            ccol1, ccol2 = st.columns(2)
            with ccol1:
                if st.button("Confirm Delete"):
                    user = st.session_state.username
                    if user in st.session_state.chat_history:
                        try:
                            st.session_state.chat_history[user].remove(st.session_state["confirm_delete_conv"])
                        except ValueError:
                            pass  # Conversation already removed.
                        save_chat_history(st.session_state.chat_history)
                    del st.session_state["confirm_delete_conv"]
                    st.sidebar.success("Conversation deleted!")
                    st.rerun()
            with ccol2:
                if st.button("Cancel"):
                    del st.session_state["confirm_delete_conv"]
                    st.sidebar.info("Deletion canceled.")
                    st.rerun()

        # --- Display chat messages with share options ---
        for idx, message in enumerate(st.session_state.messages):
            with st.chat_message(message["role"]):
                # Show role and time if available
                msg_time = message.get("time", "")
                role_title = message["role"].title()  # "User" / "Assistant"
                st.markdown(f"**`[{role_title} @ {msg_time}]`**\n\n{message['content']}")

                with st.expander("Show Copyable Text"):
                    st.code(message["content"], language="text")
                with st.expander("Share Message"):
                    st.write("### Share via WhatsApp")
                    whatsapp_contacts = load_whatsapp_contacts()
                    selected_whatsapp = st.multiselect(
                        "Choose WhatsApp Contacts:",
                        options=list(whatsapp_contacts.keys()),
                        key=f"whatsapp_select_{idx}"
                    )
                    new_whatsapp_numbers = st.text_input(
                        "Or add new WhatsApp numbers (comma-separated):",
                        key=f"whatsapp_new_{idx}"
                    )
                    st.write("### Share via Email")
                    email_contacts = load_email_contacts()
                    selected_email = st.multiselect(
                        "Choose Email Contacts:",
                        options=list(email_contacts.keys()),
                        key=f"email_select_{idx}"
                    )
                    new_email_addresses = st.text_input(
                        "Or add new Email addresses (comma-separated):",
                        key=f"email_new_{idx}"
                    )
                    if st.button("Share Message", key=f"share_button_{idx}"):
                        # Combine WhatsApp numbers from selected contacts and new entries
                        whatsapp_numbers = [whatsapp_contacts[name] for name in selected_whatsapp]
                        if new_whatsapp_numbers:
                            new_nums = [num.strip() for num in new_whatsapp_numbers.split(",") if num.strip()]
                            whatsapp_numbers.extend(new_nums)
                        whatsapp_results = {}
                        for number in whatsapp_numbers:
                            whatsapp_results[number] = send_whatsapp_message(number, f" *{current_user}*  " + message["content"])
                        
                        # Combine Email addresses from selected contacts and new entries
                        email_addresses = [email_contacts[name] for name in selected_email]
                        if new_email_addresses:
                            new_emails = [email.strip() for email in new_email_addresses.split(",") if email.strip()]
                            email_addresses.extend(new_emails)
                        email_results = {}
                        for email_addr in email_addresses:
                            email_results[email_addr] = send_email(f"NEXUS DMS Shared Chat Message from {current_user}", message["content"], email_addr)
                        
                        st.write("**WhatsApp Sharing Results:**", whatsapp_results)
                        st.write("**Email Sharing Results:**", email_results)
                
                # New expander for downloading the answer as a Word document
                with st.expander("Download as Word Document"):
                    word_buffer = create_word_doc(message["content"])
                    st.download_button(
                        label="Download Answer",
                        data=word_buffer,
                        file_name=f"Answer_{idx}.docx",  # Unique file name per message
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key=f"download_button_{idx}"  # Unique key per download button
                    )
    
                if message["role"] == "assistant":
                    with st.expander("Sources (click to view)"):
                        sources = message.get("sources", [])
                        if not sources:
                            st.write("No sources for this response.")
                        else:
                            top_k_metadata = sources[0]
                            ws_query = sources[1]

                            # Display file sources separately
                            for metadata in top_k_metadata:
                                try:
                                    st.markdown(f"**Filename:** {metadata['filename']}, **Page:** {metadata['page']}")
                                    st.code(metadata['text'], language="markdown")  # Use st.code for better formatting
                                except:
                                    st.code(json.dumps(metadata), language="markdown")  # Use st.code for better formatting

                            # Show Web Search Results separately
                            if ws_query:
                                st.markdown("Web Search Results")
                                st.code(ws_query)

                            # Button for Source Mapping
                            if st.button("Show Source Mapping", key=f"source_mapping_{int(time.time() * 1000)}"):
                                answer = message["content"]
                                with st.spinner("Mapping Source..."):
                                    final_answer = final_format(top_k_metadata, answer, ws_query)
                                    st.write(final_answer)

        # --- New User Input using text_area ---
        user_message = user_input()
        if user_message:
            ist_timezone = pytz.timezone("Asia/Kolkata")
            timestamp_now = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")

            # 1) Append the user's message to st.session_state.messages
            st.session_state.messages.append({
                "role": "user",
                "content": user_message,
                "time": timestamp_now
            })

            # Display the user message in chat
            with st.chat_message("user"):
                st.markdown(user_message)

            # [Run your retrieval code here: query_documents_with_page_range, etc.]

            # Prepare the last few messages for context.
            last_messages = st.session_state.messages[-5:] if len(st.session_state.messages) >= 5 else st.session_state.messages


            with st.spinner("Searching documents..."):
                st.markdown("**While you wait, Feel free to Refer to the Original Documents**")

                for file_key in st.session_state.selected_files:
                    # Generate the pre-signed URL for each file
                    preview_url = get_presigned_url(file_key)
                    # Create a clickable markdown link; clicking it will open the file in a new tab
                    st.markdown(f"[**{file_key}**]({preview_url})", unsafe_allow_html=True)


                top_k_metadata, answer, ws_response = query_documents_with_page_range(
                    st.session_state.selected_files, 
                    st.session_state.selected_page_ranges, 
                    user_message,
                    top_k,
                    last_messages,
                    web_search,
                    llm_model,
                    draft_mode, 
                    analyse_mode,
                    eco_mode
                )

                st.session_state.sources.append({
                    "top_k_metadata": top_k_metadata,
                    "answer": answer,
                    "websearch_metadata": ws_response
                })


            ist_timezone = pytz.timezone("Asia/Kolkata")

            # After you get the LLM's answer, you append the "assistant" message:
            assistant_answer = answer  # from your retrieval
            current_time = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
            st.session_state.messages.append({
                "role": "assistant",
                "content": assistant_answer,
                "time": current_time,
                "sources": [top_k_metadata, ws_response]
            })

            # Show the assistant response in the UI
            with st.chat_message("assistant"):
                st.markdown(assistant_answer)

            user = st.session_state.username
            if user not in st.session_state.chat_history:
                st.session_state.chat_history[user] = []

            # Check if this is a continuation of an old conversation
            original_conv_id = st.session_state.get("original_conversation_id")

            if original_conv_id:
                # This is a continuation of an old conversation
                # 1. Find the original conversation to get its label
                original_conv = None
                original_label = "Conversation"
                for conv in st.session_state.chat_history[user]:
                    if conv.get("timestamp") == original_conv_id:
                        original_conv = conv
                        original_label = conv.get("label", original_label)
                        break
            
                # 2. Create a new conversation with the old label
                new_conversation_id = str(uuid.uuid4())
                new_conversation = {
                    "conversation_id": new_conversation_id,
                    "label": original_label,
                    "timestamp": current_time,
                    "messages": st.session_state.messages,
                    "files": st.session_state.selected_files,
                    "page_ranges": st.session_state.selected_page_ranges
                }

                # 3. Remove the old conversation
                if original_conv:
                    try:
                        st.session_state.chat_history[user].remove(original_conv)
                    except ValueError:
                        pass

            # 4. Add the new conversation
                st.session_state.chat_history[user].append(new_conversation)

            # 5. Update session state
                st.session_state.current_conversation_id = new_conversation_id
                st.session_state.original_conversation_id = None  # Clear the original ID
                
            else:
                # Normal new conversation flow
                current_conversation_id = st.session_state.get("current_conversation_id")
                
                if current_conversation_id:
                    # Update existing conversation
                    for conv in st.session_state.chat_history[user]:
                        if conv.get("conversation_id") == current_conversation_id:
                            conv["messages"] = st.session_state.messages
                            conv["timestamp"] = current_time
                            conv["files"] = st.session_state.selected_files
                            conv["page_ranges"] = st.session_state.selected_page_ranges
                            break
                else:
                    # Create new conversation
                    new_conversation_id = str(uuid.uuid4())
                    st.session_state.current_conversation_id = new_conversation_id
                    
                    new_conversation = {
                        "conversation_id": new_conversation_id,
                        "label": user_message[:50],
                        "timestamp": current_time,
                        "messages": st.session_state.messages,
                        "files": st.session_state.selected_files,
                        "page_ranges": st.session_state.selected_page_ranges
                    }
                    
                    st.session_state.chat_history[user].append(new_conversation)

            # Sort and save
            st.session_state.chat_history[user] = sorted(
                st.session_state.chat_history[user],
                key=lambda x: x.get("timestamp", ""),
                reverse=True
            )

            save_chat_history(st.session_state.chat_history)
            st.rerun()

    elif option == "Query Advanced":
        st.header("Query Advanced")
        st.sidebar.header("Settings")
        llm_model = st.sidebar.selectbox("Choose Your Model", ["Skywork R1V2"])
        eco_mode = st.sidebar.toggle("Eco Mode", value=True)
        web_search = st.sidebar.toggle("Enable Web Search")
        top_k = st.sidebar.slider("Select Top-K Results", min_value=1, max_value=100, value=50, step=1)

        # File and Page Range Selection
        # available_files = list(set([metadata['filename'] for metadata in metadata_store]))
        current_user = st.session_state.get("username", "unknown")
        # Only include files where the owner is the current user or shared with the user.
        available_files = list({
            md["filename"] 
            for md in metadata_store 
            if md.get("owner") == current_user or current_user in md.get("shared_with", [])
        })


        if available_files:
            # Use multiselect and store the selection in session state.
            st.session_state.selected_files = st.multiselect(
                "Select files to include in the query:",
                available_files,
                default=st.session_state.selected_files
            )
            if len(st.session_state.selected_files) > 4:
                st.warning("For best results, select a maximum of 4 files.")
                # return

            page_ranges = get_page_range_for_files(st.session_state.selected_files)
            selected_page_ranges = {}
            # For each file, show page range inputs and store values in session state.
            for file in st.session_state.selected_files:
                min_page, max_page = page_ranges[file]
                col1, col2 = st.sidebar.columns(2)
                with col1:
                    start_page = st.number_input(
                        f"Start page for {file}",
                        min_value=min_page,
                        max_value=max_page,
                        value=page_ranges[file][0],
                        key=f"start_{file}"
                    )
                with col2:
                    end_page = st.number_input(
                        f"End page for {file}",
                        min_value=min_page,
                        max_value=max_page,
                        value=page_ranges[file][1],
                        key=f"end_{file}"
                    )
                selected_page_ranges[file] = (start_page, end_page)
            st.session_state.selected_page_ranges = selected_page_ranges
        
        
        # Input for the user to ask a question
        query = st.text_area("Ask a question about the documents (e.g., 'Compare amounts for employee benefits and management')", height=150)

        # Add a submit button for the query
        if st.button("Submit"):
            # Process the query and display the results
            if query:
                result = query_documents_viz(st.session_state.selected_files, selected_page_ranges, query, top_k, web_search, llm_model)

                # Display results or generated Python code
                if "```python" in result:
                    graphtext = result.split("```python")[0]
                    fingrphtext = remove_last_line(graphtext)
                    st.write(fingrphtext)
                    code = result.split("```python")[1].split("```")[0]
                    print(f"<<<{code}<<<")
                    try:
                        # Execute and display Plotly chart if generated
                        exec_globals = {"px": px, "go": go}
                        exec(code, exec_globals)
                        if "fig" in exec_globals:
                            st.plotly_chart(exec_globals["fig"])  # Render the Plotly chart in Streamlit
                        else:
                            st.error("No valid Plotly figure was generated in the code.")
                    except Exception as e:
                        st.error(f"Error executing code: {e}")
                else:
                    st.write(result)

    elif option == "Usage Monitoring":
        st.header("Usage Monitoring")

        # 1. Gather all usernames from your USERS dictionary
        all_usernames = list(USERS.keys())

        # 2. Add a dropdown in the sidebar for period selection
        period_options = ["Last 3 Days", "Last 7 Days", "Last 14 Days", "Last 1 Month", "Last 3 Months"]
        selected_period = st.sidebar.selectbox("Select Period", period_options, index=3)  # default: Last 1 Month
        
        # Map each period to the corresponding number of days
        period_days = {
            "Last 3 Days": 3,
            "Last 7 Days": 7,
            "Last 14 Days": 14,
            "Last 1 Month": 30,
            "Last 3 Months": 90
        }
        selected_days = period_days[selected_period]
        
        # Display the subheader with the chosen period
        st.subheader(f"Usage Monitoring - {selected_period}")

        # 3. Add a multi-select for user filtering
        selected_users = st.sidebar.multiselect(
            "Select User(s) to Display",
            options=all_usernames,
            default=all_usernames  # By default, all are selected
        )

        # 4. Load the chat history (keys are usernames)
        chat_history = load_chat_history()

        # Build a list of usage records from chat_history
        # Each record is { 'user': ..., 'timestamp': datetime_object }
        records = []
        for user, conversations in chat_history.items():
            for conv in conversations:
                timestamp_str = conv.get("timestamp")
                if timestamp_str:
                    try:
                        ts = datetime.strptime(timestamp_str, "%Y-%m-%d %H:%M:%S")
                        records.append({"user": user, "timestamp": ts})
                    except Exception as e:
                        st.warning(f"Timestamp format error for user {user}: {e}")

        # 5. Check if we have any usage records at all
        if not records:
            st.info("No usage data available.")
            return

        # Convert records to a DataFrame
        df = pd.DataFrame(records)

        # 6. Filter by the selected date range
        today = datetime.today()
        start_date = today - timedelta(days=selected_days)
        df_period = df[df["timestamp"] >= start_date]

        # 7. Filter by the selected users from the multi-select
        df_period = df_period[df_period["user"].isin(selected_users)]

        # If no data after both filters, exit
        if df_period.empty:
            st.info("No usage data found for the chosen date range and user(s).")
            return

        # --- BAR CHART: Total Queries per User ---
        user_counts = df_period.groupby("user").size().reset_index(name="queries")
        bar_fig = px.bar(
            user_counts,
            x="user",
            y="queries",
            title=f"Total Queries per User ({selected_period})",
            labels={"user": "User", "queries": "Number of Queries"}
        )
        st.plotly_chart(bar_fig, use_container_width=True)

        # --- LINE CHART: Day-wise Queries per User with Moving Average ---
        # Create a 'date' column (just the date part)
        df_period["date"] = df_period["timestamp"].dt.date

        # Count queries per user per day
        daily_counts = df_period.groupby(["user", "date"]).size().reset_index(name="queries")

        # Create a complete date range for the selected period
        date_range = pd.date_range(start=start_date.date(), end=today.date())
        all_users_in_data = daily_counts["user"].unique()
        complete_data = []

        # Build a day-by-day dataset with possible missing dates filled in
        for user in all_users_in_data:
            user_df = daily_counts[daily_counts["user"] == user].copy()
            user_df.set_index("date", inplace=True)
            # Reindex to fill missing days with 0
            user_df = user_df.reindex(date_range, fill_value=0)
            user_df = user_df.rename_axis("date").reset_index()
            user_df["user"] = user
            
            # Calculate a moving average based on 'selected_days'
            user_df["moving_avg"] = user_df["queries"].rolling(window=selected_days, min_periods=1).mean()
            complete_data.append(user_df)

        # Combine each user's timeseries data
        daily_all = pd.concat(complete_data, ignore_index=True)

        # Create the line graph showing daily queries
        line_fig = px.line(
            daily_all,
            x="date",
            y="queries",
            color="user",
            title=f"Daily Queries per User ({selected_period})",
            labels={"date": "Date", "queries": "Number of Queries"}
        )

        # Add a separate trace for the rolling average of each user
        for user in all_users_in_data:
            user_data = daily_all[daily_all["user"] == user]
            line_fig.add_trace(
                go.Scatter(
                    x=user_data["date"],
                    y=user_data["moving_avg"],
                    mode="lines",
                    name=f"{user} - {selected_period} MA"
                )
            )

        st.plotly_chart(line_fig, use_container_width=True)

    else:
        st.warning("No files available in the index. Please upload Documents to populate the index.")

if __name__ == "__main__":
    main()